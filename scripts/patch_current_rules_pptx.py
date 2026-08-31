"""Patch the current rewritten rules PPTX in place. Does not rebuild the booklet."""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Mm, Pt

ROOT = Path(__file__).resolve().parents[1]
ART = ROOT / "assets" / "rules"
PPTX = ROOT / "output" / "svet-moy-zerkalce-rules.pptx"
EMU_PER_MM = 914400 / 25.4
BODY_PT = 11.0
BOTTOM_MARGIN_MM = 12.0


def _mm(emu: int) -> float:
    return emu / EMU_PER_MM


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def _first_run_size(shape) -> float | None:
    if not shape.has_text_frame:
        return None
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                return r.font.size.pt
    return None


def _is_heading(shape) -> bool:
    if not shape.has_text_frame:
        return False
    text = _shape_text(shape)
    if not text:
        return False
    size = _first_run_size(shape)
    if size is not None and size >= 17:
        return True
    first = text.split("\n", 1)[0].strip()
    bold = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.bold:
                bold = True
                break
    return bool(bold and first == first.upper() and len(first) < 40)


def _is_caption(shape) -> bool:
    if not shape.has_text_frame:
        return False
    text = _shape_text(shape).lower()
    if not text:
        return False
    keys = ("шт.", "карты букв:", "карты голосования")
    return any(key in text for key in keys) and _mm(shape.width) < 80


def _is_komplekt_caption(shape, slide_idx: int) -> bool:
    if slide_idx != 0 or not _is_caption(shape):
        return False
    return "шт." in _shape_text(shape).lower()


def _in_compose_section(shape, slide_idx: int) -> bool:
    """Body/example/picture/caption of «Составление ответа», not the heading."""
    if slide_idx != 1:
        return False
    text = _shape_text(shape)
    if text.startswith("СОСТАВЛЕНИЕ"):
        return False
    top = _mm(shape.top)
    return 18.0 <= top < 60.5


def replace_blip(shape, image_path: Path) -> None:
    _part, rid = shape.part.get_or_add_image_part(str(image_path))
    blip = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip")
    if blip is None:
        raise RuntimeError("picture has no blip")
    blip.set(qn("r:embed"), rid)


def set_picture_size(shape, width_mm: float, height_mm: float) -> None:
    shape.width = Mm(width_mm)
    shape.height = Mm(height_mm)


def patch_fonts(prs: Presentation) -> None:
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _is_heading(shape):
                continue
            if _is_komplekt_caption(shape, si):
                continue
            old = _first_run_size(shape)
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(BODY_PT)
            if old is not None and old < BODY_PT - 0.1:
                ratio = BODY_PT / old
                shape.height = Emu(int(shape.height * ratio))


def patch_alignment(prs: Presentation) -> None:
    for si, slide in enumerate(prs.slides):
        texts = [
            sh
            for sh in slide.shapes
            if sh.has_text_frame and _shape_text(sh) and _mm(sh.width) >= 70
        ]
        if not texts:
            continue
        left_target = min(sh.left for sh in texts if _is_heading(sh) or _mm(sh.width) >= 150)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _is_caption(shape):
                continue
            if _in_compose_section(shape, si):
                continue
            if _mm(shape.width) < 70:
                continue
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
            if _mm(shape.width) >= 150:
                shape.left = left_target
            elif abs(_mm(shape.left) - _mm(left_target)) < 8:
                shape.left = left_target


def _group_bounds(shapes: list) -> tuple[int, int]:
    top = min(sh.top for sh in shapes)
    bottom = max(sh.top + sh.height for sh in shapes)
    return top, bottom


def stack_groups(groups: list[list], slide_height: int) -> None:
    groups = [g for g in groups if g]
    if not groups:
        return
    heights = []
    rels = []
    for g in groups:
        top, bottom = _group_bounds(g)
        heights.append(bottom - top)
        rels.append([(sh, sh.top - top, sh.left) for sh in g])
    start = _group_bounds(groups[0])[0]
    budget = slide_height - Mm(BOTTOM_MARGIN_MM) - start - sum(heights)
    gap = int(budget / (len(groups) - 1)) if len(groups) > 1 else 0
    min_gap = int(Mm(2.5))
    if gap < min_gap:
        gap = min_gap
    y = start
    for g, h, rel in zip(groups, heights, rels):
        for sh, rel_top, _left in rel:
            sh.top = y + rel_top
        y += h + gap


def patch_pictures(prs: Presentation) -> None:
    answer = ART / "card-letter-answer.png"
    setup = ART / "setup-table.png"
    aw, ah = Image.open(answer).size
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            w_mm, h_mm = _mm(shape.width), _mm(shape.height)
            # Setup illustration: wide landscape under подготовка
            if si == 0 and w_mm > 85 and h_mm > 50:
                replace_blip(shape, setup)
            # Compose-answer illustration: currently letter duplex, left of body
            if si == 1 and 50 < w_mm < 70 and 25 < h_mm < 40:
                replace_blip(shape, answer)
                new_h = w_mm * (ah / aw)
                set_picture_size(shape, w_mm, new_h)
                pic_left, pic_w = shape.left, shape.width
                pic_bottom = shape.top + shape.height
                for cap in slide.shapes:
                    if not _is_caption(cap):
                        continue
                    if "карты букв:" in _shape_text(cap).lower():
                        cap.top = pic_bottom + Mm(0.4)
                        cap.left = pic_left
                        cap.width = pic_w


def slide_groups(slide, si: int) -> list[list]:
    shapes = list(slide.shapes)
    if si == 0:
        return [
            [shapes[0], shapes[1], shapes[2]],
            [shapes[3], shapes[4], shapes[5], shapes[6], shapes[7], shapes[8], shapes[9]],
            [shapes[10], shapes[11], shapes[12]],
            [shapes[13], shapes[14]],
            [shapes[15], shapes[16], shapes[17]],
        ]
    if si == 1:
        return [
            [shapes[0], shapes[1], shapes[2], shapes[3]],
            [shapes[5], shapes[6], shapes[7], shapes[8], shapes[9]],
            [shapes[10], shapes[11]],
            [shapes[12], shapes[13]],
            [shapes[14], shapes[15], shapes[18], shapes[19], shapes[20], shapes[21]],
            [shapes[16], shapes[17]],
        ]
    return []


def patch() -> Path:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from generate_rules_art import build as build_art

    build_art()
    prs = Presentation(str(PPTX))
    patch_pictures(prs)
    patch_fonts(prs)
    patch_alignment(prs)
    for si, slide in enumerate(prs.slides):
        groups = slide_groups(slide, si)
        if groups:
            stack_groups(groups, prs.slide_height)
    prs.save(str(PPTX))
    return PPTX


if __name__ == "__main__":
    path = patch()
    print(f"Wrote {path}")
