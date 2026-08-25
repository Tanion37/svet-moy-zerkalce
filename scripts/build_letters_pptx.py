"""A4-portrait PnP for letter cards (Russian frequency deck).

Design card 57×44.1 mm. Grid 3×6 = 18 per page; total tiles scaled to 72.
Follows TanionAgentSetting/prototype-presentation.md (labels, duplex, word-fit).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Mm, Pt

try:
    from PIL import ImageFont
except ImportError:  # pragma: no cover
    ImageFont = None  # type: ignore

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"

FREQ_WEIGHTS: dict[str, int] = {
    "о": 1097,
    "е": 848,
    "а": 806,
    "и": 737,
    "н": 670,
    "т": 631,
    "с": 547,
    "р": 473,
    "в": 454,
    "л": 440,
    "к": 349,
    "м": 321,
    "д": 298,
    "п": 281,
    "у": 262,
    "я": 201,
    "ы": 190,
    "ь": 174,
    "г": 169,
    "з": 165,
    "б": 159,
    "ч": 147,
    "й": 121,
    "х": 97,
    "ж": 94,
    "ш": 73,
    "ю": 64,
    "ц": 48,
    "щ": 36,
    "э": 32,
    "ф": 26,
    "ё": 20,
    "ъ": 2,
}

DESIGN_CARD_W_MM = 57.0
DESIGN_CARD_H_MM = 44.1
COLS = 3
ROWS = 6
PER_SLIDE = COLS * ROWS

SLIDE_W_MM = 210.0
SLIDE_H_MM = 297.0
PRINTER_MARGIN_MM = 5.0
LABEL_BAND_MM = 6.0

BORDER_EMU = 9525
LINE_MM = BORDER_EMU / 36000.0
SAFE_INSET_MM = PRINTER_MARGIN_MM + LINE_MM / 2
USABLE_W_MM = SLIDE_W_MM - 2 * SAFE_INSET_MM
USABLE_H_MM = SLIDE_H_MM - 2 * SAFE_INSET_MM - LABEL_BAND_MM

DESIGN_GRID_W = COLS * DESIGN_CARD_W_MM
DESIGN_GRID_H = ROWS * DESIGN_CARD_H_MM
SCALE = min(USABLE_W_MM / DESIGN_GRID_W, USABLE_H_MM / DESIGN_GRID_H)
CARD_W_MM = DESIGN_CARD_W_MM * SCALE
CARD_H_MM = DESIGN_CARD_H_MM * SCALE
CONTENT_W = COLS * CARD_W_MM
CONTENT_H = ROWS * CARD_H_MM
MARGIN_X_MM = (SLIDE_W_MM - CONTENT_W) / 2
MARGIN_Y_MM = SAFE_INSET_MM + LABEL_BAND_MM + (
    (SLIDE_H_MM - SAFE_INSET_MM - LABEL_BAND_MM - SAFE_INSET_MM - CONTENT_H) / 2
)

BORDER_COLOR = RGBColor(0xC0, 0xC0, 0xC0)
INK = RGBColor(0x1A, 0x1A, 0x1A)
LABEL_COLOR = RGBColor(0x55, 0x55, 0x55)
LETTER_PT = 28.0
MIN_PT = 8.0
LABEL_PT = 9.0
PAD_X_MM = 1.5
PAD_Y_MM = 1.0

TOTAL_LETTERS = ((70 + PER_SLIDE - 1) // PER_SLIDE) * PER_SLIDE  # 72


def build_weighted_deck(weights: dict[str, int], total: int) -> list[str]:
    letters = sorted(weights.keys(), key=lambda ch: (ch.replace("ё", "е\uffff"), ch))
    if total < len(letters):
        raise ValueError(f"total {total} < alphabet {len(letters)}")

    assigned = {k: 1 for k in letters}
    remain = total - len(letters)

    entries = [(k, weights[k]) for k in letters if weights[k] > 0]
    weight_sum = sum(w for _, w in entries)
    raw = []
    for k, w in entries:
        exact = (w / weight_sum) * remain
        raw.append({"k": k, "floor": int(exact), "frac": exact - int(exact)})
    used = sum(x["floor"] for x in raw)
    raw.sort(key=lambda x: x["frac"], reverse=True)
    i = 0
    while used < remain:
        raw[i % len(raw)]["floor"] += 1
        used += 1
        i += 1
    for x in raw:
        assigned[x["k"]] += x["floor"]

    out: list[str] = []
    for k in letters:
        out.extend([k] * assigned[k])
    return out


def _no_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _add_run(
    paragraph,
    text: str,
    size: float,
    *,
    bold: bool = True,
    color: RGBColor = INK,
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _word_width_pt(word: str, size_pt: float, *, bold: bool) -> float:
    if ImageFont is not None:
        try:
            name = "arialbd.ttf" if bold else "arial.ttf"
            font = ImageFont.truetype(name, int(round(size_pt)))
            return float(font.getlength(word))
        except OSError:
            pass
    return len(word) * size_pt * (0.72 if bold else 0.58)


def fit_font_for_words(
    text: str,
    max_width_mm: float,
    max_pt: float,
    *,
    bold: bool = True,
    min_pt: float = MIN_PT,
) -> float:
    words = [w for w in text.replace("\n", " ").split(" ") if w]
    if not words:
        return max_pt
    max_width_pt = max_width_mm * 72.0 / 25.4
    size = max_pt
    while size > min_pt + 1e-6:
        widest = max(_word_width_pt(w, size, bold=bold) for w in words)
        if widest <= max_width_pt:
            return round(size, 1)
        size -= 0.5
    return min_pt


def add_sheet_label(slide, text: str) -> None:
    box = slide.shapes.add_textbox(
        Mm(SAFE_INSET_MM),
        Mm(SAFE_INSET_MM),
        Mm(SLIDE_W_MM - 2 * SAFE_INSET_MM),
        Mm(LABEL_BAND_MM - 0.5),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Mm(0)
    tf.margin_right = Mm(0)
    tf.margin_top = Mm(0)
    tf.margin_bottom = Mm(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, text, LABEL_PT, bold=False, color=LABEL_COLOR)


def _cut_hline(slide, left_mm: float, top_mm: float, width_mm: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(left_mm),
        Mm(top_mm - LINE_MM / 2),
        Mm(width_mm),
        Mm(LINE_MM),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER_COLOR
    shape.line.fill.background()
    _no_shadow(shape)


def _cut_vline(slide, left_mm: float, top_mm: float, height_mm: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(left_mm - LINE_MM / 2),
        Mm(top_mm),
        Mm(LINE_MM),
        Mm(height_mm),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER_COLOR
    shape.line.fill.background()
    _no_shadow(shape)


def add_cut_grid(slide, left0: float, top0: float) -> None:
    occupied = {(r, c) for r in range(ROWS) for c in range(COLS)}
    for row_edge in range(ROWS + 1):
        for col in range(COLS):
            above = (row_edge - 1, col) in occupied if row_edge > 0 else False
            below = (row_edge, col) in occupied if row_edge < ROWS else False
            if above or below:
                _cut_hline(
                    slide,
                    left0 + col * CARD_W_MM,
                    top0 + row_edge * CARD_H_MM,
                    CARD_W_MM,
                )
    for col_edge in range(COLS + 1):
        for row in range(ROWS):
            left = (row, col_edge - 1) in occupied if col_edge > 0 else False
            right = (row, col_edge) in occupied if col_edge < COLS else False
            if left or right:
                _cut_vline(
                    slide,
                    left0 + col_edge * CARD_W_MM,
                    top0 + row * CARD_H_MM,
                    CARD_H_MM,
                )


def add_letter_card(slide, left_mm: float, top_mm: float, letter: str) -> None:
    text_w = CARD_W_MM - 2 * PAD_X_MM * SCALE
    size = fit_font_for_words(letter.upper(), text_w, LETTER_PT, bold=True)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(left_mm),
        Mm(top_mm),
        Mm(CARD_W_MM),
        Mm(CARD_H_MM),
    )
    shape.fill.background()
    shape.line.fill.background()
    _no_shadow(shape)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Mm(PAD_X_MM * SCALE)
    tf.margin_right = Mm(PAD_X_MM * SCALE)
    tf.margin_top = Mm(PAD_Y_MM * SCALE)
    tf.margin_bottom = Mm(PAD_Y_MM * SCALE)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.space_before = Pt(0)
    p0.space_after = Pt(0)
    _add_run(p0, letter.upper(), size, bold=True)


def add_face_sheet(slide, left0: float, top0: float, chunk: list[str]) -> None:
    assert len(chunk) == PER_SLIDE
    for idx, letter in enumerate(chunk):
        col = idx % COLS
        row = idx // COLS
        add_letter_card(
            slide,
            left0 + col * CARD_W_MM,
            top0 + row * CARD_H_MM,
            letter,
        )
    add_cut_grid(slide, left0, top0)
    add_sheet_label(slide, "Карты букв · ЛИЦО")


def add_back_sheet(slide, left0: float, top0: float, chunk: list[str]) -> None:
    """Empty joker backs, column-mirrored; no cut lines."""
    assert len(chunk) == PER_SLIDE
    for row in range(ROWS):
        for face_col in range(COLS):
            back_col = COLS - 1 - face_col
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Mm(left0 + back_col * CARD_W_MM),
                Mm(top0 + row * CARD_H_MM),
                Mm(CARD_W_MM),
                Mm(CARD_H_MM),
            )
            shape.fill.background()
            shape.line.fill.background()
            _no_shadow(shape)
    add_sheet_label(slide, "Карты букв · ОБОРОТ")


def _chunks(items: list[str], size: int) -> list[list[str]]:
    return [items[i : i + size] for i in range(0, len(items), size)]


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build() -> Path:
    assert TOTAL_LETTERS % PER_SLIDE == 0
    letters = build_weighted_deck(FREQ_WEIGHTS, TOTAL_LETTERS)
    assert len(letters) == TOTAL_LETTERS

    prs = Presentation()
    prs.slide_width = Mm(SLIDE_W_MM)
    prs.slide_height = Mm(SLIDE_H_MM)
    left0, top0 = MARGIN_X_MM, MARGIN_Y_MM

    for chunk in _chunks(letters, PER_SLIDE):
        add_face_sheet(_blank_slide(prs), left0, top0, chunk)
        add_back_sheet(_blank_slide(prs), left0, top0, chunk)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "svet-moy-zerkalce-letters.pptx"
    prs.save(str(out))
    return out, letters


if __name__ == "__main__":
    path, letters = build()
    from collections import Counter

    counts = Counter(letters)
    print(f"Wrote {path}")
    print(f"Slide {SLIDE_W_MM:.2f}x{SLIDE_H_MM:.2f} mm (A4 portrait)")
    print(
        f"Grid {COLS}x{ROWS}={PER_SLIDE}; total {TOTAL_LETTERS} "
        f"({TOTAL_LETTERS // PER_SLIDE} face/back pairs)"
    )
    print(
        f"Cards {CARD_W_MM:.3f}x{CARD_H_MM:.3f} mm "
        f"(design {DESIGN_CARD_W_MM}x{DESIGN_CARD_H_MM}, scale {SCALE:.4f})"
    )
    print("Counts:", dict(sorted(counts.items(), key=lambda x: (-x[1], x[0]))))
