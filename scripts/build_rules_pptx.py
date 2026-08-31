"""Rules PPTX in Hobby World (МХ) booklet tone.

A4 portrait. Component art from assets/rules (PnP look, not named placeholders).
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Mm, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
ART = ROOT / "assets" / "rules"

SLIDE_W = 210.0
SLIDE_H = 297.0
MARGIN = 12.0
CONTENT_W = SLIDE_W - 2 * MARGIN
MAX_GAP = 29.0  # ~5 lines

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0x90, 0x90, 0x90)
ACCENT = RGBColor(0x2E, 0x5A, 0x88)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CONTACTS = (
    "t.me/tanion, vk.ru/tanion, yury@nxt.ru, +79213189331 "
    "(не стесняйтесь звонить / писать при возникновении вопросов)"
)
RULES_VERSION = "3.7"


def _no_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _run(p, text, size, *, bold=False, italic=False, color=INK):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def _textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(width), Mm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Mm(0)
    tf.margin_right = Mm(0)
    tf.margin_top = Mm(0)
    tf.margin_bottom = Mm(0)
    return tf


def add_title(slide, text: str, top: float) -> float:
    tf = _textbox(slide, MARGIN, top, CONTENT_W, 14)
    p = tf.paragraphs[0]
    _run(p, text, 18, bold=True)
    return top + 15


def add_heading(slide, text: str, top: float) -> float:
    tf = _textbox(slide, MARGIN, top, CONTENT_W, 8)
    p = tf.paragraphs[0]
    _run(p, text.upper(), 12, bold=True, color=ACCENT)
    return top + 8


def add_body(slide, lines: list[str], top: float, *, size: float = 11, italic: bool = False) -> float:
    h = max(7.0, len(lines) * (size * 0.42 + 2.0) + 2)
    tf = _textbox(slide, MARGIN, top, CONTENT_W, h)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        _run(p, line, size, italic=italic)
    return top + h + 1


def add_picture(slide, name: str, left: float, top: float, width: float) -> float:
    path = ART / name
    pic = slide.shapes.add_picture(str(path), Mm(left), Mm(top), Mm(width))
    return pic.height.mm


class Book:
    def __init__(self, prs: Presentation):
        self.prs = prs
        self.slide = None
        self.y = MARGIN

    def leftover(self) -> float:
        return SLIDE_H - MARGIN - self.y

    def new_slide(self) -> None:
        self.slide = _blank(self.prs)
        self.y = MARGIN

    def need(self, h: float) -> None:
        if self.slide is None:
            self.new_slide()
            return
        if self.leftover() >= h:
            return
        self.new_slide()

    def heading(self, text: str) -> None:
        self.need(10)
        self.y = add_heading(self.slide, text, self.y)

    def title(self, text: str) -> None:
        self.need(14)
        self.y = add_title(self.slide, text, self.y)

    def body(self, lines: list[str], *, size: float = 11, italic: bool = False) -> None:
        h = max(7.0, len(lines) * (size * 0.42 + 2.0) + 3)
        self.need(h)
        self.y = add_body(self.slide, lines, self.y, size=size, italic=italic)

    def _img_size(self, name: str, width: float) -> tuple[float, float]:
        im = Image.open(ART / name)
        return width, width * im.height / im.width

    def picture_row(
        self,
        items: list[tuple[str, str]],
        *,
        card_w: float = 42.0,
        gap: float = 4.0,
        caption_h: float = 8.0,
    ) -> None:
        n = len(items)
        sizes = [self._img_size(name, card_w) for name, _ in items]
        max_h = max(h for _, h in sizes)
        total = n * card_w + (n - 1) * gap
        left0 = MARGIN + max(0.0, (CONTENT_W - total) / 2)
        block = max_h + caption_h + 2
        self.need(block)
        for i, ((name, caption), (_w, card_h)) in enumerate(zip(items, sizes)):
            left = left0 + i * (card_w + gap)
            top = self.y + (max_h - card_h)
            add_picture(self.slide, name, left, top, card_w)
            tf = _textbox(self.slide, left, self.y + max_h + 0.5, card_w, caption_h)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _run(p, caption, 8, color=MUTED)
        self.y += block

    def picture_request_stack(
        self,
        face_name: str,
        back_name: str,
        face_cap: str,
        back_cap: str,
        *,
        card_w: float = 70.0,
        gap: float = 3.0,
        caption_h: float = 8.0,
    ) -> None:
        """Face above back – the request-card layout used everywhere in the booklet."""
        _w, card_h = self._img_size(face_name, card_w)
        block = 2 * (card_h + caption_h) + gap + 2
        self.need(block)
        left = MARGIN + (CONTENT_W - card_w) / 2
        y = self.y
        add_picture(self.slide, face_name, left, y, card_w)
        tf = _textbox(self.slide, left, y + card_h + 0.4, card_w, caption_h)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, face_cap, 8, color=MUTED)
        y += card_h + caption_h + gap
        add_picture(self.slide, back_name, left, y, card_w)
        tf = _textbox(self.slide, left, y + card_h + 0.4, card_w, caption_h)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, back_cap, 8, color=MUTED)
        self.y += block

    def component_row(self) -> None:
        """Комплектация: ряд картинок, подписи только под ними."""
        items = [
            (
                "card-request-duplex.png",
                "карты запросов – 14 шт. (лицевая сторона и рубашка)",
            ),
            (
                "card-letter-duplex.png",
                "карты букв – 72 шт. (рубашка – джокер)",
            ),
            (
                "card-vote-stack.png",
                "карты «Ты восхитителен» – 6 шт. (по 1 каждого цвета)",
            ),
        ]
        gap = 4.0
        caption_h = 14.0
        n = len(items)
        aspects = []
        for name, _ in items:
            im = Image.open(ART / name)
            aspects.append(im.height / im.width)
        # Fit row into CONTENT_W; cap height so the row stays compact.
        max_h = 42.0
        widths = [min(max_h / a, 70.0) for a in aspects]
        total = sum(widths) + (n - 1) * gap
        if total > CONTENT_W:
            scale = CONTENT_W / total
            widths = [w * scale for w in widths]
        heights = [w * a for w, a in zip(widths, aspects)]
        row_h = max(heights)
        block = row_h + caption_h + 2
        self.need(block)
        x = MARGIN + (CONTENT_W - (sum(widths) + (n - 1) * gap)) / 2
        y0 = self.y
        for (name, caption), cw, ch in zip(items, widths, heights):
            top = y0 + (row_h - ch)
            add_picture(self.slide, name, x, top, cw)
            tf = _textbox(self.slide, x, y0 + row_h + 0.4, cw, caption_h)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _run(p, caption, 8, color=MUTED)
            x += cw + gap
        self.y += block

    def side_by_side(
        self,
        lines: list[str],
        image_name: str,
        *,
        image_on_right: bool = True,
        img_w: float = 58.0,
        size: float = 11,
        italic: bool = False,
        caption: str | None = None,
        caption_h: float = 8.0,
    ) -> None:
        """Narrow illustration: text on one half, image on the opposite half."""
        _w, img_h = self._img_size(image_name, img_w)
        cap = caption_h if caption else 0.0
        gap = 5.0
        text_w = CONTENT_W - img_w - gap
        text_h = max(7.0, len(lines) * (size * 0.42 + 2.2) + 2)
        block = max(text_h, img_h + cap) + 2
        self.need(block)
        if image_on_right:
            text_left = MARGIN
            img_left = MARGIN + text_w + gap
        else:
            img_left = MARGIN
            text_left = MARGIN + img_w + gap
        tf = _textbox(self.slide, text_left, self.y, text_w, max(text_h, img_h))
        first = True
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            _run(p, line, size, italic=italic)
        img_top = self.y
        add_picture(self.slide, image_name, img_left, img_top, img_w)
        if caption:
            ctf = _textbox(
                self.slide, img_left, img_top + img_h + 0.3, img_w, caption_h
            )
            p = ctf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _run(p, caption, 8, color=MUTED)
        self.y += block

    def setup_image(self, max_h: float | None = None) -> None:
        name = "setup-table.png"
        native_w, native_h = 200.0, 130.0
        w = CONTENT_W
        h = w * (native_h / native_w)
        if max_h is not None and h > max_h:
            h = max_h
            w = h * (native_w / native_h)
        self.need(min(h, self.leftover() if self.leftover() > MAX_GAP else h))
        if self.leftover() < h:
            scale = self.leftover() / h
            w *= scale
            h *= scale
        left = MARGIN + (CONTENT_W - w) / 2
        add_picture(self.slide, name, left, self.y, w)
        self.y += h + 2

    def step_row(self, steps: list[tuple[str, str]]) -> None:
        n = len(steps)
        box_w = (CONTENT_W - (n - 1) * 4) / n
        h = 22.0
        self.need(h + 2)
        x = MARGIN
        for i, (num, title) in enumerate(steps):
            sh = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Mm(x), Mm(self.y), Mm(box_w), Mm(h)
            )
            sh.fill.solid()
            sh.fill.fore_color.rgb = LIGHT
            sh.line.color.rgb = LINE
            sh.line.width = Pt(0.75)
            _no_shadow(sh)
            tf = sh.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _run(p, f"{num}. {title}", 9, bold=True)
            x += box_w + 4
        self.y += h + 3


def build_book(book: Book) -> None:
    book.title("Свет мой зеркальце. Правила игры")
    book.body(
        [
            "3–6 игроков  ·  от 16 лет  ·  около 30 минут",
            f"Версия правил {RULES_VERSION}",
        ],
        size=12,
    )
    book.body(
        [
            "Вы – гремлины службы поддержки волшебного сервиса «Свет мой зеркальце скажи». Клиенты из почти современного мира задают через зеркало странные вопросы, а вы на лету собираете ответы из карт букв, которые берёте из общей кучи в центре стола. Чем смешнее и точнее реакция зала – тем выше шанс забрать карту запроса себе. Побеждает тот, у кого к концу колоды запросов больше всего таких карт.",
        ],
        size=11,
        italic=True,
    )
    book.heading("Комплектация")
    book.component_row()
    book.body(
        [
            "Элементы, которых нет в PnP, но которые нужны для партии: поверхность стола, достаточная чтобы раскидать карты букв и кидать карты «Ты восхитителен» к игрокам.",
        ],
        size=10,
    )

    book.heading("Подготовка к игре")
    book.body(
        [
            "1. Перемешайте колоду карт запросов и положите её рубашками вверх. Верхняя карта колоды запросов показывает окончание будущего запроса на своей рубашке.",
            "2. Каждый игрок выбирает цвет и берёт карту «Ты восхитителен» своего цвета. Лишние карты «Ты восхитителен» уберите в коробку.",
            "3. Все карты букв раскиньте в случайном порядке лицом вверх в центре стола. Руки игроков в начале раунда пустые.",
        ]
    )
    book.setup_image()

    book.heading("Цель игры")
    book.body(
        [
            "Набирайте карты запросов как победные очки. Когда колода запросов больше не позволяет начать раунд, побеждает игрок с наибольшим числом таких карт.",
        ]
    )
    book.heading("Как устроена карта запроса")
    book.side_by_side(
        [
            "У каждой карты запроса две стороны.",
            "1. Лицевая сторона – начало фразы-запроса.",
            "2. Рубашка – окончание фразы-запроса.",
            "Рамка зеркала на всех картах запросов одинаковая: при любой перетасовке лицевая сторона одной карты запроса и рубашка следующей сверху карты в колоде всегда складываются в одно изображение зеркала с целым текстом запроса.",
            "Читаете вслух, например: «Расскажи чем испугать всех соседей сверху».",
        ],
        "card-request-pair.png",
        image_on_right=True,
        img_w=52.0,
        size=10,
        caption="лицевая сторона сверху, рубашка снизу",
    )

    book.heading("Ход игры")
    book.body(
        [
            "Игра идёт раундами. Все действия в раунде выполняются одновременно.",
        ]
    )
    book.step_row(
        [
            ("1", "Открытие запроса"),
            ("2", "Составление ответа"),
            ("3", "Голосование"),
            ("4", "Возврат букв"),
        ]
    )
    book.body(
        [
            "1. Снимите верхнюю карту колоды запросов, переверните её на лицевую сторону и положите рядом с колодой так, чтобы вместе с новой верхней картой колоды запросов (с её рубашкой) получилось общее зеркало.",
            "Прочитайте вслух текст на лицевой стороне открытой карты запроса и на рубашке верхней карты запроса.",
            "Если карты кончились, наступает конец игры.",
            "2. Игроки одновременно берут карты букв из центра стола одной рукой по одной карте. Второй рукой можно держать взятое, но не копаться в куче.",
            "Берите карты букв, пока вам хватит на ответ. Лишние карты букв остаются в центре стола.",
            "Любую карту буквы можно сыграть рубашкой вверх как джокер (любая буква). Замену вслух не называйте.",
            "3. Голосование – см. ниже.",
            "4. Все карты букв возвращаются в центр стола, перемешиваются и снова раскидываются лицом вверх.",
        ],
        size=10,
    )

    book.heading("Составление ответа")
    book.side_by_side(
        [
            "Ответ – слово или фраза на усмотрение игрока в пределах взятых карт букв. Допустимы несуществующие и неполные слова.",
            "Несыгранные взятые карты букв остаются у игрока до конца раунда.",
            "Игрок выкладывает свой ответ перед собой на стол, когда готов.",
        ],
        "card-letter-duplex.png",
        image_on_right=False,
        img_w=62.0,
        size=10,
        caption="карты букв: лицо и рубашка-джокер",
    )

    book.body(
        [
            "Пример: запрос «Чем соблазнить милф». Настя выложила «БЕЗ ТРУСОВ» (позже получит 3 голоса), Артём – «В ПОДЪЕЗДЕ» (1 голос), Кирилл – «МАМИН ДРУГ» (0 голосов).",
        ],
        size=11,
        italic=True,
    )

    book.heading("Голосование")
    book.side_by_side(
        [
            "Каждый игрок берёт в руку карту голосования. Когда все определились с выбором, то по команде одновременно кидают на понравившееся слово (если очень хочется – можно в составившего его игрока). Если неясно, к кому упала – кинувший поясняет.",
            "Себе карту голосования оставлять нельзя: кидайте только к чужому ответу.",
            "Если игрок не успел кинуть карту голосования вместе с остальными, его карта не учитывается, а каждый другой игрок получает по 2 голоса.",
            "Победитель раунда – игрок, перед которым оказалось больше всего карт голосования. Он берёт карту запроса, лежащую лицевой стороной вверх, и кладёт её перед собой как победное очко.",
            "Ничья по голосам. Если лидеров несколько, все они объявляются победителями раунда. Один из них берёт лицевую карту запроса, остальные по очереди (по часовой стрелке от обладателя лицевой карты запроса или по договорённости) берут по одной карте с верха колоды запросов как победные очки. Часть будущих запросов при этом выбывает – так и задумано.",
            "После определения победителя или победителей верните карты «Ты восхитителен» владельцам в руку – они снова понадобятся в следующем раунде.",
        ],
        "card-vote-stack.png",
        image_on_right=True,
        img_w=48.0,
        size=10,
        caption="карты голосования «Ты восхитителен»",
    )

    book.heading("Конец игры")
    book.body(
        [
            "Игра заканчивается, когда колода запросов больше не позволяет начать раунд: карты кончились или осталась одна карта запроса без пары для зеркала.",
            "1. Каждый игрок считает карты запросов, лежащие перед ним как победные очки.",
            "2. Побеждает игрок с наибольшим числом таких карт.",
            "При ничьей объявите совместную победу.",
        ]
    )
    book.heading("Частые вопросы")
    book.body(
        [
            "? Можно ли кинуть карту голосования себе?  Нет. Только к чужому ответу.",
            "? Карта голосования упала между двумя игроками – чей это голос?  Кинувший поясняет.",
            "? Можно ли копаться в разложенных картах букв?  Нет. Берите только сверху или с края кучи, одной рукой по одной карте. Второй рукой можно держать взятое, но не копаться в куче.",
            "? Когда переставать брать карты букв?  Когда вам хватит карт букв на ваш ответ.",
            "? Сколько джокеров можно сыграть?  Сколько угодно. Зал реже голосует за нечитаемый ответ.",
        ],
        size=10,
    )
    book.heading("Создатели")
    book.body(
        [
            "Автор: Юрий Ямщиков",
            "Редакция / развитие: [плейсхолдер]",
            f"Версия правил {RULES_VERSION}",
        ],
        size=10,
    )
    book.heading("Памятка раунда")
    book.body(
        [
            "1. Открыть карту запроса: снять верхнюю карту колоды запросов, перевернуть на лицевую сторону, прочитать текст на лицевой стороне открытой карты запроса и на рубашке верхней карты запроса.",
            "2. Взять карты букв из центра стола (одна рука, по одной карте) и собрать ответ.",
            "3. Кинуть карту голосования на понравившееся слово.",
            "4. Все карты букв вернуть в центр стола, перемешать и раскидать лицом вверх.",
        ]
    )
    book.body(
        [
            "После раунда все карты букв снова лежат в центре стола лицом вверх.",
        ]
    )
    book.setup_image()
    book.heading("Контакты")
    book.body(
        [
            f"Контакты автора: {CONTACTS}",
        ],
        size=10,
    )


def build() -> Path:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from generate_rules_art import build as build_art

    build_art()

    prs = Presentation()
    prs.slide_width = Mm(SLIDE_W)
    prs.slide_height = Mm(SLIDE_H)

    book = Book(prs)
    build_book(book)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "svet-moy-zerkalce-rules.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
