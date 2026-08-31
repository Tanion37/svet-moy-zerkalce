"""A4-landscape PnP: request faces; request backs; «Ты восхитителен» as its own file.

КОРНИ 15: рубашки компонентов — отдельным файлом.
Follows TanionAgentSetting/prototype-presentation.md:
- word-fit font shrink; sheet labels; no back page for single-sided;
- back sheet still column-mirrored for duplex with the face sheet.
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Mm, Pt

try:
    from PIL import ImageFont
except ImportError:  # pragma: no cover
    ImageFont = None  # type: ignore

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
ASSETS = ROOT / "assets"
FACE_FRAME = ASSETS / "request-face-frame.png"
BACK_FRAME = ASSETS / "request-back-frame.png"

DESIGN_CARD_W_MM = 88.0
DESIGN_CARD_H_MM = 63.0
COLS = 3
ROWS = 3
PER_SLIDE = COLS * ROWS

SLIDE_W_MM = 297.0
SLIDE_H_MM = 210.0
PRINTER_MARGIN_MM = 5.0
LABEL_BAND_MM = 6.0

BORDER_EMU = 9525
LINE_MM = BORDER_EMU / 36000.0
SAFE_INSET_MM = PRINTER_MARGIN_MM + LINE_MM / 2
USABLE_W_MM = SLIDE_W_MM - 2 * SAFE_INSET_MM
# Leave a top band inside the safe area for the sheet label.
USABLE_H_MM = SLIDE_H_MM - 2 * SAFE_INSET_MM - LABEL_BAND_MM

DESIGN_GRID_W = COLS * DESIGN_CARD_W_MM
DESIGN_GRID_H = ROWS * DESIGN_CARD_H_MM
SCALE = min(1.0, USABLE_W_MM / DESIGN_GRID_W, USABLE_H_MM / DESIGN_GRID_H)
CARD_W_MM = DESIGN_CARD_W_MM * SCALE
CARD_H_MM = DESIGN_CARD_H_MM * SCALE
CONTENT_W = COLS * CARD_W_MM
CONTENT_H = ROWS * CARD_H_MM
MARGIN_X_MM = (SLIDE_W_MM - CONTENT_W) / 2
MARGIN_Y_MM = SAFE_INSET_MM + LABEL_BAND_MM + (
    (SLIDE_H_MM - SAFE_INSET_MM - LABEL_BAND_MM - SAFE_INSET_MM - CONTENT_H) / 2
)

BORDER_COLOR = RGBColor(0xC0, 0xC0, 0xC0)
INK = RGBColor(0x1A, 0x1A, 0x1A)
LABEL_COLOR = RGBColor(0x1A, 0x1A, 0x1A)

MAIN_PT = 20.0
MIN_PT = 8.0
PLAYER_LABEL_PT = 10.0
LABEL_PT = 11.0
PAD_X_MM = 2.4
PAD_Y_MM = 2.0
# Keep text inside the open area of the arch / U-frame (not on the rails).
FRAME_SIDE_FR = 0.11
FRAME_END_FR = 0.36
FRAME_EDGE_MM = 3.5

FACES: list[str] = [
    "Заклинание превращения в одного из",
    "Расскажи чем испугать всех",
    "Подбери прозвище для кого-то из",
    "Придумай страшное проклятие для",
    "Комплимент, который вгонит в краску",
    "Подскажи как отшить",
    "Почему стоит избегать",
    "Какая польза от",
    "Назови тайную слабость",
    "Какое наказание ждёт в аду для",
    "Чем оскорбить любого из",
    "Какой подлянки ждать от",
    "Лучший подарок для",
    "Чем соблазнить",
]
# Extra wrap: keep MAIN_PT, narrower column, N lines.
FACE_WRAP_LINES: dict[str, int] = {
    "Подбери прозвище для кого-то из": 3,
}

BACKS: list[str] = [
    "конченных интровертов",
    "здесь присутствующих",
    "ухоженных дровосеков",
    "твоих бывших",
    "соседей сверху",
    "любителей настолок",
    "самокатчиков",
    "бабок у подъезда",
    "чемпионов френдзоны",
    "пузатых скуфов",
    "милф",
    "душнил",
    "нейросетей",
    "очкариков",
]

PLAYER_COLORS: list[RGBColor] = [
    RGBColor(0xC6, 0x28, 0x28),
    RGBColor(0x15, 0x65, 0xC0),
    RGBColor(0x2E, 0x7D, 0x32),
    RGBColor(0xEF, 0x6C, 0x00),
    RGBColor(0x6A, 0x1B, 0x9A),
    RGBColor(0x00, 0x83, 0x8F),
]


def _no_shadow(shape) -> None:
    """Drop theme effectRef (shadow) — inherit=False alone is not enough."""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    el = shape._element
    style = el.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
    )
    if style is not None:
        el.remove(style)
    sp_pr = getattr(el, "spPr", None)
    if sp_pr is None:
        sp_pr = el.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr"
        )
    if sp_pr is not None:
        for child in list(sp_pr):
            if child.tag.endswith("}effectLst"):
                sp_pr.remove(child)


def set_picture_fill(shape, image_path: Path) -> None:
    """Blip-fill the card shape so text stays on the same component (§7)."""
    _image_part, rId = shape.part.get_or_add_image_part(str(image_path))
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag.endswith(
            ("}solidFill", "}noFill", "}gradFill", "}blipFill", "}pattFill", "}grpFill")
        ):
            sp_pr.remove(child)
    blip_fill = etree.Element(qn("a:blipFill"))
    blip = etree.SubElement(blip_fill, qn("a:blip"))
    blip.set(qn("r:embed"), rId)
    stretch = etree.SubElement(blip_fill, qn("a:stretch"))
    etree.SubElement(stretch, qn("a:fillRect"))
    xfrm = sp_pr.find(qn("a:xfrm"))
    if xfrm is not None:
        xfrm.addnext(blip_fill)
    else:
        sp_pr.insert(0, blip_fill)


def _add_run(
    paragraph,
    text: str,
    size: float,
    *,
    bold: bool = False,
    color: RGBColor = INK,
    font_name: str = "Arial",
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _word_width_pt(word: str, size_pt: float, *, bold: bool) -> float:
    if ImageFont is not None:
        try:
            name = "arialbd.ttf" if bold else "arial.ttf"
            font = ImageFont.truetype(name, int(round(size_pt)))
            return float(font.getlength(word))
        except OSError:
            pass
    # Fallback: conservative bold Cyrillic estimate
    return len(word) * size_pt * (0.72 if bold else 0.58)


def fit_font_for_words(
    text: str,
    max_width_mm: float,
    max_pt: float,
    *,
    bold: bool = True,
    min_pt: float = MIN_PT,
) -> float:
    """Shrink until every word fits on one line within max_width_mm."""
    words = [w for w in text.replace("\n", " ").split(" ") if w]
    if not words:
        return max_pt
    max_width_pt = max_width_mm * 72.0 / 25.4
    size = max_pt
    while size > min_pt + 1e-6:
        widest = max(_word_width_pt(w, size, bold=bold) for w in words)
        if widest <= max_width_pt:
            return round(size, 1)
        size -= 0.5
    return min_pt


def _wrap_words(
    text: str, width_mm: float, size_pt: float, *, bold: bool = True
) -> list[str]:
    words = [w for w in text.replace("\n", " ").split(" ") if w]
    if not words:
        return []
    space = _word_width_pt(" ", size_pt, bold=bold)
    width_px = width_mm * 72.0 / 25.4

    def phrase_w(ws: list[str]) -> float:
        return sum(_word_width_pt(w, size_pt, bold=bold) for w in ws) + space * (
            len(ws) - 1
        )

    lines: list[str] = []
    cur: list[str] = []
    for word in words:
        trial = cur + [word]
        if not cur or phrase_w(trial) <= width_px:
            cur = trial
        else:
            lines.append(" ".join(cur))
            cur = [word]
    if cur:
        lines.append(" ".join(cur))
    return lines


def extra_side_for_n_lines(
    text: str,
    width_mm: float,
    n_lines: int,
    size_pt: float,
    *,
    bold: bool = True,
) -> float:
    """Extra mm on each side so greedy wrap is n_lines at this size."""
    words = [w for w in text.replace("\n", " ").split(" ") if w]
    if not words:
        return 0.0
    longest_mm = max(_word_width_pt(w, size_pt, bold=bold) for w in words) * 25.4 / 72.0
    max_extra_total = max(0.0, width_mm - longest_mm)
    hits: list[float] = []
    extra = 0.0
    step = 0.1
    while extra <= max_extra_total + 1e-9:
        wrapped = _wrap_words(text, width_mm - extra, size_pt, bold=bold)
        if len(wrapped) == n_lines:
            hits.append(extra)
        extra += step
    if not hits:
        return 0.0
    # Clearly narrower than the 2-line edge, still above the 4-line cliff.
    extra_total = hits[0] + 0.6 * (hits[-1] - hits[0])
    return extra_total / 2.0


def add_sheet_label(slide, text: str) -> None:
    """Label inside printer-safe area, above the card grid."""
    box = slide.shapes.add_textbox(
        Mm(SAFE_INSET_MM),
        Mm(SAFE_INSET_MM),
        Mm(SLIDE_W_MM - 2 * SAFE_INSET_MM),
        Mm(LABEL_BAND_MM - 0.5),
    )
    _no_shadow(box)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Mm(0)
    tf.margin_right = Mm(0)
    tf.margin_top = Mm(0)
    tf.margin_bottom = Mm(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, text, LABEL_PT, bold=True, color=LABEL_COLOR)


def _cut_hline(slide, left_mm: float, top_mm: float, width_mm: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(left_mm),
        Mm(top_mm - LINE_MM / 2),
        Mm(width_mm),
        Mm(LINE_MM),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER_COLOR
    shape.line.fill.background()
    _no_shadow(shape)


def _cut_vline(slide, left_mm: float, top_mm: float, height_mm: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(left_mm - LINE_MM / 2),
        Mm(top_mm),
        Mm(LINE_MM),
        Mm(height_mm),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER_COLOR
    shape.line.fill.background()
    _no_shadow(shape)


def add_cut_grid(
    slide, left0: float, top0: float, occupied: set[tuple[int, int]] | None = None
) -> None:
    if occupied is None:
        occupied = {(r, c) for r in range(ROWS) for c in range(COLS)}

    for row_edge in range(ROWS + 1):
        for col in range(COLS):
            above = (row_edge - 1, col) in occupied if row_edge > 0 else False
            below = (row_edge, col) in occupied if row_edge < ROWS else False
            if above or below:
                _cut_hline(
                    slide,
                    left0 + col * CARD_W_MM,
                    top0 + row_edge * CARD_H_MM,
                    CARD_W_MM,
                )
    for col_edge in range(COLS + 1):
        for row in range(ROWS):
            left = (row, col_edge - 1) in occupied if col_edge > 0 else False
            right = (row, col_edge) in occupied if col_edge < COLS else False
            if left or right:
                _cut_vline(
                    slide,
                    left0 + col_edge * CARD_W_MM,
                    top0 + row * CARD_H_MM,
                    CARD_H_MM,
                )


def _frame_margins(*, face: bool) -> tuple[float, float, float, float]:
    side = CARD_W_MM * FRAME_SIDE_FR
    deep = CARD_H_MM * FRAME_END_FR
    edge = FRAME_EDGE_MM * SCALE
    if face:
        return side, side, deep, edge
    return side, side, edge, deep


def _card_shape(
    slide,
    left_mm: float,
    top_mm: float,
    *,
    anchor: MSO_ANCHOR,
    art: Path | None = None,
    margins_mm: tuple[float, float, float, float] | None = None,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(left_mm),
        Mm(top_mm),
        Mm(CARD_W_MM),
        Mm(CARD_H_MM),
    )
    shape.line.fill.background()
    _no_shadow(shape)
    if art is None:
        shape.fill.background()
    else:
        set_picture_fill(shape, art)

    if margins_mm is None:
        pad_x = PAD_X_MM * SCALE
        pad_y = PAD_Y_MM * SCALE
        left_m = right_m = pad_x
        top_m = bottom_m = pad_y
    else:
        left_m, right_m, top_m, bottom_m = margins_mm

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Mm(left_m)
    tf.margin_right = Mm(right_m)
    tf.margin_top = Mm(top_m)
    tf.margin_bottom = Mm(bottom_m)
    return tf


def _text_width_mm(margins_mm: tuple[float, float, float, float] | None = None) -> float:
    if margins_mm is None:
        return CARD_W_MM - 2 * PAD_X_MM * SCALE
    left_m, right_m, _top, _bottom = margins_mm
    return CARD_W_MM - left_m - right_m


def add_text_card(
    slide,
    left_mm: float,
    top_mm: float,
    text: str,
    *,
    anchor: MSO_ANCHOR,
    art: Path | None = None,
    margins_mm: tuple[float, float, float, float] | None = None,
    bold: bool = True,
    color: RGBColor = INK,
    size: float | None = None,
) -> None:
    if size is None:
        size = fit_font_for_words(text, _text_width_mm(margins_mm), MAIN_PT, bold=bold)
    tf = _card_shape(
        slide, left_mm, top_mm, anchor=anchor, art=art, margins_mm=margins_mm
    )
    lines = text.split("\n") or [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        _add_run(p, line, size, bold=bold, color=color)


def add_ty_voskhititelen_card(
    slide, left_mm: float, top_mm: float, player: int, color: RGBColor
) -> None:
    title = "Ты восхитителен"
    size = fit_font_for_words(title, _text_width_mm(), MAIN_PT, bold=True)
    tf = _card_shape(slide, left_mm, top_mm, anchor=MSO_ANCHOR.MIDDLE)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.space_before = Pt(0)
    p0.space_after = Pt(0)
    _add_run(p0, title, size, bold=True, color=color)

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(6)
    p1.space_after = Pt(0)
    _add_run(p1, f"Игрок {player}", PLAYER_LABEL_PT, bold=False, color=color)


def _cell(left0: float, top0: float, idx: int) -> tuple[float, float, int, int]:
    col = idx % COLS
    row = idx // COLS
    return left0 + col * CARD_W_MM, top0 + row * CARD_H_MM, row, col


def _chunks(items: list[str], size: int) -> list[list[str]]:
    return [items[i : i + size] for i in range(0, len(items), size)]


def add_request_face_sheet(slide, left0: float, top0: float, chunk: list[str]) -> None:
    margins = _frame_margins(face=True)
    occupied: set[tuple[int, int]] = set()
    for idx, text in enumerate(chunk):
        left, top, row, col = _cell(left0, top0, idx)
        occupied.add((row, col))
        card_margins = margins
        card_text = text
        size: float | None = None
        n_lines = FACE_WRAP_LINES.get(text)
        if n_lines:
            extra = extra_side_for_n_lines(
                text, _text_width_mm(margins), n_lines, MAIN_PT
            )
            left_m, right_m, top_m, bottom_m = margins
            card_margins = (left_m + extra, right_m + extra, top_m, bottom_m)
            card_text = "\n".join(
                _wrap_words(text, _text_width_mm(card_margins), MAIN_PT)
            )
            size = MAIN_PT
        add_text_card(
            slide,
            left,
            top,
            card_text,
            anchor=MSO_ANCHOR.BOTTOM,
            art=FACE_FRAME,
            margins_mm=card_margins,
            size=size,
        )
    add_cut_grid(slide, left0, top0, occupied)
    add_sheet_label(slide, "Карты запросов · ЛИЦО")


def add_request_back_sheet(slide, left0: float, top0: float, chunk: list[str]) -> None:
    """Column-mirrored backs; no cut lines."""
    margins = _frame_margins(face=False)
    n = len(chunk)
    rows_used = (n + COLS - 1) // COLS
    for row in range(rows_used):
        row_count = min(COLS, n - row * COLS)
        for face_col in range(row_count):
            back_col = COLS - 1 - face_col
            face_idx = row * COLS + face_col
            add_text_card(
                slide,
                left0 + back_col * CARD_W_MM,
                top0 + row * CARD_H_MM,
                chunk[face_idx],
                anchor=MSO_ANCHOR.TOP,
                art=BACK_FRAME,
                margins_mm=margins,
            )
    add_sheet_label(slide, "Карты запросов · ОБОРОТ")


def add_vote_sheet(slide, left0: float, top0: float) -> None:
    """Single-sided «Ты восхитителен» — no back slide."""
    occupied: set[tuple[int, int]] = set()
    for i, color in enumerate(PLAYER_COLORS):
        left, top, row, col = _cell(left0, top0, i)
        occupied.add((row, col))
        add_ty_voskhititelen_card(slide, left, top, i + 1, color)
    add_cut_grid(slide, left0, top0, occupied)
    add_sheet_label(slide, "Карты «Ты восхитителен»")


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Mm(SLIDE_W_MM)
    prs.slide_height = Mm(SLIDE_H_MM)
    return prs


def build() -> tuple[Path, Path, Path]:
    assert CARD_W_MM > CARD_H_MM
    assert len(FACES) == len(BACKS) == 14
    if not FACE_FRAME.is_file() or not BACK_FRAME.is_file():
        raise FileNotFoundError(f"Missing frame art: {FACE_FRAME} / {BACK_FRAME}")

    left0, top0 = MARGIN_X_MM, MARGIN_Y_MM
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    faces = _new_prs()
    backs = _new_prs()
    for face_chunk, back_chunk in zip(_chunks(FACES, PER_SLIDE), _chunks(BACKS, PER_SLIDE)):
        add_request_face_sheet(_blank_slide(faces), left0, top0, face_chunk)
        add_request_back_sheet(_blank_slide(backs), left0, top0, back_chunk)

    faces_out = OUT_DIR / "svet-moy-zerkalce-cards.pptx"
    faces.save(str(faces_out))

    backs_out = OUT_DIR / "svet-moy-zerkalce-card-backs.pptx"
    backs.save(str(backs_out))

    votes = _new_prs()
    add_vote_sheet(_blank_slide(votes), left0, top0)
    votes_out = OUT_DIR / "svet-moy-zerkalce-votes.pptx"
    votes.save(str(votes_out))

    return faces_out, backs_out, votes_out


if __name__ == "__main__":
    faces_path, backs_path, votes_path = build()
    face_margins = _frame_margins(face=True)
    sample = "здесь присутствующих"
    fitted = fit_font_for_words(sample, _text_width_mm(face_margins), MAIN_PT)
    print(f"Wrote {faces_path}")
    print(f"Wrote {backs_path}")
    print(f"Wrote {votes_path}")
    print(f"Slide {SLIDE_W_MM:.2f}x{SLIDE_H_MM:.2f} mm (exact A4 landscape)")
    print(
        f"Margins >= {PRINTER_MARGIN_MM} mm "
        f"(actual x={MARGIN_X_MM:.3f}, y={MARGIN_Y_MM:.3f})"
    )
    print(
        f"Cards {CARD_W_MM:.3f}x{CARD_H_MM:.3f} mm "
        f"(design {DESIGN_CARD_W_MM}x{DESIGN_CARD_H_MM}, scale {SCALE:.4f})"
    )
    nick = "Подбери прозвище для кого-то из"
    nick_extra = extra_side_for_n_lines(nick, _text_width_mm(face_margins), 3, MAIN_PT)
    nick_w = _text_width_mm(face_margins) - 2 * nick_extra
    nick_lines = _wrap_words(nick, nick_w, MAIN_PT)
    print(f"Word-fit sample «{sample}»: {fitted} pt (max {MAIN_PT})")
    print(
        f"«{nick}»: {len(nick_lines)} lines at {MAIN_PT} pt, "
        f"width {nick_w:.2f} mm (was {_text_width_mm(face_margins):.2f}), "
        f"{nick_lines}"
    )
    n_face_sheets = (len(FACES) + PER_SLIDE - 1) // PER_SLIDE
    print(
        f"Faces/backs {len(FACES)} on {n_face_sheets} duplex sheet(s); "
        f"votes {len(PLAYER_COLORS)} in a separate face-only file"
    )
